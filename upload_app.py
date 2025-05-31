import re
import os
import uuid
import json
import numpy as np 
from flask import Flask, request, jsonify, render_template, url_for, send_from_directory, session
from flask_cors import CORS
from werkzeug.utils import secure_filename
from gtts import gTTS
from transformers import pipeline
import pdf2image
import pytesseract
import docx
import pptx
import google.generativeai as genai
from dotenv import load_dotenv
from PIL import Image
from sentence_transformers import SentenceTransformer, util
import pickle
import base64
from io import BytesIO
import os.path

from document_utils import store_document_data, get_document_data, clear_document_data

app = Flask(__name__)
CORS(app)

app.secret_key = os.urandom(24)

UPLOAD_FOLDER = 'uploads'
AUDIO_FOLDER = os.path.join('static', 'audio_files')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(AUDIO_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['AUDIO_FOLDER'] = AUDIO_FOLDER

load_dotenv()
api_key = os.getenv('GOOGLE_API_KEY')
app.config['GOOGLE_API_KEY'] = api_key

if not api_key:
    raise ValueError("GOOGLE_API_KEY not found in environment variables")

genai.configure(api_key=api_key)

qa_model = pipeline("question-answering")
similarity_model = SentenceTransformer('all-MiniLM-L6-v2')

def clean_text(text):
    cleaned_text = re.sub(r'\*+|[_~`^]', '', text)
    cleaned_text = re.sub(r'[^\w\s.,!?]', '', cleaned_text)
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text).strip()
    return cleaned_text

def extract_text(file_path):
    lower_path = file_path.lower()
    print("[DEBUG] Extracting text from:", lower_path)

    if lower_path.endswith('.pdf'):
        images = pdf2image.convert_from_path(file_path)
        text = ""
        for image in images:
            text += pytesseract.image_to_string(image) + "\n"
        return text
    elif lower_path.endswith('.docx'):
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    elif lower_path.endswith('.pptx'):
        presentation = pptx.Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif lower_path.endswith(('.jpg', '.jpeg', '.png', '.bmp', '.tiff')):
        image = Image.open(file_path)
        return pytesseract.image_to_string(image)
    else:
        raise ValueError("Unsupported file type")

@app.route('/talk', methods=['POST'])
def talk():
    data = request.json
    text = data.get('text', '')

    if not text:
        return jsonify({'error': 'No text provided'}), 400

    filename = f"{uuid.uuid4().hex}.mp3"
    filepath = os.path.join(AUDIO_FOLDER, filename)

    tts = gTTS(text)
    tts.save(filepath)

    return jsonify({
        'audio_url': f'/static/audio/{filename}',
        'blendData': text
    })

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/get_session_data', methods=['GET'])
def get_session_data():
    return jsonify({
        'extracted_text': session.get('extracted_content', ''),
        'explanation': session.get('explanation', ''),
        'audio_file': session.get('audio_file', ''),
        'quiz': session.get('quiz', ''),
        'has_processed_file': 'extracted_content' in session
    })

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file part'}), 400
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No selected file'}), 400

        if 'user_id' not in session:
            session['user_id'] = str(uuid.uuid4())
        
        user_id = session['user_id']
        
        clear_document_data(user_id)
        
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)

        print("[DEBUG] File saved:", file_path)

        extracted_text = extract_text(file_path)
        print("[DEBUG] Extracted text preview:", repr(extracted_text[:200]))

        model = genai.GenerativeModel('models/gemini-2.0-flash')
        explanation = model.generate_content(
            f"Explain this content simply and clearly: {extracted_text}"
        ).text.strip()

        cleaned_explanation = clean_text(explanation)

        unique_filename = f"{uuid.uuid4().hex}.mp3"
        audio_path = os.path.join(app.config['AUDIO_FOLDER'], unique_filename)
        tts = gTTS(text=cleaned_explanation, lang='en')
        tts.save(audio_path)

        quiz_prompt = f"""Content: {extracted_text}

Generate a quiz in strict JSON:
[
    {{
        "question": "Question text",
        "option": ["Option1", "Option2", "Option3", "Option4"],
        "answer": "Correct answer",
        "explanation": "Explanation text"
    }},
    ...
]
Rules:
- Avoid repeating or paraphrasing.
- Return only valid JSON.
"""
        quiz_response = model.generate_content(quiz_prompt).text.strip()

        if quiz_response.startswith("```json"):
            quiz_response = quiz_response.lstrip("```json").rstrip("```").strip()
        elif quiz_response.startswith("```"):
            quiz_response = quiz_response.lstrip("```").rstrip("```").strip()

        doc_data = {
            "extracted_content": extracted_text,
            "explanation": cleaned_explanation,
            "quiz": quiz_response
        }
        
        store_document_data(user_id, doc_data)
        
        session['has_document'] = True
        session['audio_file'] = url_for('static', filename=f'audio_files/{unique_filename}')
        
        print("[DEBUG] Document data saved to filesystem for user:", user_id)

        return jsonify({
            'extracted_text': extracted_text,
            'explanation': cleaned_explanation,
            'audio_file': url_for('static', filename=f'audio_files/{unique_filename}'),
            'quiz': quiz_response,
            'avatar_video': ''
        })

    except Exception as e:
        print("UPLOAD ERROR:", str(e))
        return jsonify({'error': f"Server Error: {str(e)}"}), 500

@app.route('/generate_quiz', methods=['POST'])
def generate_quiz():
    try:
        data = request.get_json()
        extracted_text = data.get("extracted_text", "")
        if not extracted_text:
            return jsonify({"error": "No extracted text provided"}), 400

        previous_questions = data.get("previous_questions", [])
        
        difficulty_level = data.get("difficulty", "mixed") 
        taxonomy_level = data.get("taxonomy", "mixed")  
        question_format = data.get("format", "mixed") 
        
        similarity_model = SentenceTransformer('all-MiniLM-L6-v2')
        
        if previous_questions:
            previous_embeddings = similarity_model.encode(previous_questions)
        else:
            previous_embeddings = np.array([])
        
        model = genai.GenerativeModel('models/gemini-2.0-flash')

        taxonomy_explanations = {
            "knowledge": "recall facts, terms, basic concepts",
            "comprehension": "demonstrate understanding of facts and ideas",
            "application": "apply knowledge to new situations",
            "analysis": "examine information and break it down",
            "evaluation": "present and defend opinions by making judgments",
            "creation": "create new ideas or ways of viewing things"
        }
        
        format_instructions = {
            "mcq": "Create only multiple-choice questions with 4 options each",
            "true_false": "Create only true/false questions",
            "short_answer": "Create open-ended questions requiring short answers"
        }
        
        taxonomy_instruction = ""
        if taxonomy_level != "mixed" and taxonomy_level in taxonomy_explanations:
            taxonomy_instruction = f"- Questions should focus on {taxonomy_level} level ({taxonomy_explanations[taxonomy_level]})"
        elif taxonomy_level == "mixed":
            taxonomy_instruction = "- Include questions covering various levels of Bloom's taxonomy (knowledge, comprehension, application, analysis, evaluation, creation)"
            
        difficulty_instruction = ""
        if difficulty_level != "mixed":
            difficulty_instruction = f"- All questions should be {difficulty_level} difficulty"
        else:
            difficulty_instruction = "- Include a mix of easy, medium, and hard difficulty questions"
            
        format_instruction = ""
        if question_format != "mixed" and question_format in format_instructions:
            format_instruction = f"- {format_instructions[question_format]}"
        elif question_format == "mixed":
            format_instruction = "- Include a mix of multiple-choice, true/false, and short-answer questions"

        prompt = f"""Content: {extracted_text}

Generate a quiz in strict JSON format:
[
    {{
        "question": "Question text",
        "option": ["Option1", "Option2", "Option3", "Option4"],
        "answer": "Correct answer",
        "explanation": "Explanation of correct answer",
        "difficulty": "easy|medium|hard",
        "taxonomy_level": "knowledge|comprehension|application|analysis|evaluation|creation",
        "format": "mcq|true_false|short_answer"
    }},
    ...
]
Important:
{difficulty_instruction}
{taxonomy_instruction}
{format_instruction}
- Create diverse questions covering different concepts from the content
- Ensure questions test different aspects of understanding
- Create questions that differ substantially from each other
- For true/false questions, options should be ["True", "False"]
- For short-answer questions, include a list of acceptable answers in the "option" field
- Return valid JSON only.
"""

        quiz_response = model.generate_content(prompt).text.strip()

        if quiz_response.startswith("```json"):
            quiz_response = quiz_response.lstrip("```json").rstrip("```").strip()
        elif quiz_response.startswith("```"):
            quiz_response = quiz_response.lstrip("```").rstrip("```").strip()

        try:
            quiz_data = json.loads(quiz_response)
        except json.JSONDecodeError as e:
            return jsonify({"error": f"Gemini returned invalid JSON: {str(e)}", "raw_response": quiz_response}), 500

        unique_quiz = []
        seen_questions = set()
        
        def is_semantically_similar(new_question, prev_embeddings, threshold=0.85):
            if not isinstance(prev_embeddings, np.ndarray) or prev_embeddings.size == 0:  
                return False
                
            new_embedding = similarity_model.encode([new_question])[0]
            
            similarities = util.pytorch_cos_sim(new_embedding, prev_embeddings)[0]
            
            return any(sim.item() > threshold for sim in similarities)

        for item in quiz_data:
            question_text = item.get("question", "").strip()
            
            if not question_text or question_text.lower() in seen_questions:
                continue
                
            if isinstance(previous_embeddings, np.ndarray) and previous_embeddings.size > 0 and is_semantically_similar(question_text, previous_embeddings):
                print(f"Skipping semantically similar question: {question_text}")
                continue
                
            seen_questions.add(question_text.lower())
            unique_quiz.append(item)
            
            new_embedding = similarity_model.encode([question_text])
            if isinstance(previous_embeddings, np.ndarray) and previous_embeddings.size > 0:
                previous_embeddings = np.vstack([previous_embeddings, new_embedding])
            else:
                previous_embeddings = new_embedding

        if not unique_quiz:
            return jsonify({"error": "No new quiz questions could be generated."}), 400

        return jsonify({"quiz": json.dumps(unique_quiz)})

    except Exception as e:
        print("QUIZ GENERATION ERROR:", str(e))
        return jsonify({'error': f"Quiz Generation Failed: {str(e)}"}), 500
    
@app.route('/uploads/<path:filename>', methods=['GET'])
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

from avatar_flask_routes import register_avatar_routes
register_avatar_routes(app)

if __name__ == '__main__':
    app.run(debug=True, port=5000)
